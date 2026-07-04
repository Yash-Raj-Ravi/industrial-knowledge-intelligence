from pptx import Presentation

def parse_pptx(file_path:str) -> str:
    presentation = Presentation(file_path)
    slides_text=[]
    for slide in presentation.slides:
        for shape in slide .shapes:
            if shape.has_text_frame: # has_text_frame prevents errors when a shape doesn't contain text.
                text = shape.text
                if text.strip():
                    slides_text.append(text)
    return "\n".join(slides_text)